#!/usr/bin/env python3
import argparse
from pathlib import Path
import csv
from pptx import Presentation

def extract_text_from_shape(shape):
    """Extract text from a shape if it has text."""
    if hasattr(shape, "text"):
        return shape.text.strip()
    return ""

def extract_text_from_slide(slide):
    """Extract all text from a slide by iterating through its shapes."""
    texts = []
    for shape in slide.shapes:
        text = extract_text_from_shape(shape)
        if text:
            texts.append(text)
    return " ".join(texts)

def process_presentation(pptx_path):
    """Process PowerPoint file and save text content to CSV."""
    prs = Presentation(pptx_path)
    output_path = pptx_path.with_suffix('.csv')
    
    with open(output_path, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f, quoting=csv.QUOTE_ALL)
        writer.writerow(['slide_number', 'text_content'])
        
        for idx, slide in enumerate(prs.slides, 1):
            text_content = extract_text_from_slide(slide)
            writer.writerow([str(idx), text_content])
    
    return output_path

def main():
    parser = argparse.ArgumentParser(description='Extract text from PowerPoint slides')
    parser.add_argument('input_file', type=str, help='Path to PowerPoint file')
    
    args = parser.parse_args()
    input_path = Path(args.input_file)
    
    if not input_path.exists():
        print(f"Error: Input file '{input_path}' does not exist")
        return 1
    
    if not input_path.suffix.lower() in ['.pptx', '.ppt']:
        print(f"Error: Input file must be a PowerPoint file (.pptx or .ppt)")
        return 1
        
    try:
        output_path = process_presentation(input_path)
        print(f"Successfully extracted text to {output_path}")
        return 0
    except Exception as e:
        print(f"Error processing presentation: {e}")
        return 1

if __name__ == '__main__':
    exit(main()) 